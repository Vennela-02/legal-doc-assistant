import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from PyPDF2 import PdfReader
from typing import List, Tuple
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image
from dotenv import load_dotenv
load_dotenv()

TESSERACT_CMD = os.getenv("TESSERACT_CMD", "tesseract")
POPPLER_PATH = os.getenv("POPPLER_PATH", "/usr/bin")

def extract_text_from_file(filename: str, content: bytes) -> List[Tuple[str, int, str]]:
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == ".txt":
        return extract_text_from_txt(content, filename)

    elif ext == ".pdf":
        return extract_text_from_pdf(content, filename)

    elif ext == ".docx":
        return extract_text_from_docx(content, filename)

    elif ext == ".pptx":
        return extract_text_from_pptx(content, filename)

    else:
        raise ValueError(f"Unsupported file format: {ext}")



# Apply Tesseract path
pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")
POPPLER_PATH = os.getenv("POPPLER_PATH")

def extract_text_with_ocr_from_bytes(content: bytes) -> List[Tuple[str, int]]:
    """Run OCR on every page and return (text, page_number)."""
    
    if POPPLER_PATH:  
        pages = convert_from_bytes(content, dpi=300, poppler_path=POPPLER_PATH)
    else:  
        pages = convert_from_bytes(content, dpi=300)

    ocr_pages = []
    for i, page_img in enumerate(pages):
        text = pytesseract.image_to_string(page_img)
        ocr_pages.append((text.strip(), i + 1))

    return ocr_pages


def extract_text_from_pdf(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    """Extract text from PDF using PyMuPDF + OCR hybrid approach (merged per page)."""
    pages = []

    #  Extract digital text with PyMuPDF
    with fitz.open(stream=content, filetype="pdf") as doc:
        digital_text = [(page.get_text().strip(), i + 1, source) for i, page in enumerate(doc)]

    # Extract OCR text
    ocr_pages = extract_text_with_ocr_from_bytes(content)

    # Merge per page
    for (dig_text, page_num, src), (ocr_text, ocr_page_num) in zip(digital_text, ocr_pages):
        assert page_num == ocr_page_num  # safety check
        # Combine but avoid duplication
        if dig_text and ocr_text:
            merged = dig_text + "\n\n[OCR Extracted]\n" + ocr_text
        elif dig_text:
            merged = dig_text
        else:
            merged = ocr_text
        pages.append((merged.strip(), page_num, src))

    print(f"✅ Hybrid extraction complete for {len(pages)} pages.")
    return pages

def extract_text_from_docx(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    from io import BytesIO
    doc = docx.Document(BytesIO(content))
    full_text = "\n".join([para.text for para in doc.paragraphs])
    return [(full_text, 1, source)]  # Treat as one page

def extract_text_from_pptx(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    from io import BytesIO
    prs = Presentation(BytesIO(content))
    pages = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        pages.append((slide_text.strip(), i + 1, source))
    return pages

def extract_text_from_txt(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    text = content.decode("utf-8", errors="ignore")
    return [(text, 1, source)]  # Treat as one page
